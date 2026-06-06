# -*- coding: utf-8 -*-
"""
build_pptx.py — 把 REPORT.md 的 15 頁藍圖生成 SudokuCoach.pptx
  · 每頁:標題 + 重點條列 + 對應的圖(測試/訓練/demo 都嵌進去)
  · 🎤口條放進每頁的『簡報者備忘稿』(presenter notes)
圖來源:experiments/figures/(先跑過 experiments/*.py 與 report/make_report_figs.py)
執行:python report/build_pptx.py  →  產出 report/SudokuCoach.pptx
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(ROOT, "experiments", "figures")
FONT = "Microsoft JhengHei"
NAVY = RGBColor(0x1A, 0x3A, 0x5A)
GREY = RGBColor(0x44, 0x44, 0x44)
ACCENT = RGBColor(0x6A, 0x1B, 0x9A)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def fig(name):
    p = os.path.join(FIG, name)
    return p if os.path.exists(p) else None


# (標題, 副標, [重點], [圖檔], 口條)
SLIDES = [
    ("數獨教練 SudokuCoach",
     "從感知到推理:以數獨為例的神經符號系統實作 — A Neuro-Symbolic Approach",
     [], [], "標題頁。副標放英文,格局直接出來。"),

    ("破題:「丟給 ChatGPT 不就好了?」",
     "白盒 vs 機率",
     ["LLM 是機率模型:解釋是『事後編的』,可能跟計算無關、甚至自信地出錯",
      "本系統推理是白盒:理由是『算出來的』,每步可驗證、保證正確",
      "原則:有精確解的問題,別用機率的鎚子敲"],
     [],
     "把全場最想問的問題先端上桌,框出白盒 vs 機率的主軸。"),

    ("系統總覽:兩個半腦",
     "神經符號架構(背景)",
     ["感知 CNN:會『看』,把照片變成 81 個數字 — 黑盒,說不出理由",
      "推理引擎:會『想』,用數獨規則一步步推 — 白盒,每步有理由",
      "本報告主角是前半段那顆 CNN:它學到什麼、何時會錯、錯了怎麼辦"],
     [],
     "兩者互補,叫 Neuro-Symbolic;但這份報告的主角是 CNN。"),

    ("感知 I:前處理 + 為什麼不用 MNIST",
     "灰階→二值化→找盤面→透視校正→切 81 格",
     ["流程:adaptiveThreshold → 最大凸四邊形(epsilon 掃描)→ 透視校正 → 切格",
      "不用 MNIST:數獨是『印刷體』、MNIST 是『手寫』,分布不同",
      "改用字型渲染合成資料(digit_dataset.py)+ 隨機擾動",
      "⚠️ 自備:放自己拍的照片『前處理 before/after』截圖"],
     [],
     "訓練/測試分布要一致,模型才有效——所以資料是用字型印出來的。"),

    ("感知 II:CNN 架構",
     "Conv→Pool→Conv→Pool→Dense→Softmax",
     ["形狀變化:28→26→13→11→5→1600→128→10",
      "Conv 學濾波器、Pool 給平移容忍、Dense 做分類",
      "前半段抽特徵、後半段做決定;參數大多在 Dense(Conv 省又會看)",
      "想更直觀:現場開 learn/lab.py 手算一次卷積"],
     [],
     "卷積在學一組濾波器、池化給平移容忍、全連接做分類。"),

    ("CNN 學到什麼:濾波器 + 特徵圖",
     "特徵是『長出來的』,不是寫死的",
     ["第一層 32 個濾波器:自學出邊緣/筆畫偵測器",
      "特徵圖:一個真實數字觸發了哪些通道"],
     ["filters.png", "feature_maps.png"],
     "第一層自己長出了邊緣偵測器——沒有人教它,是從資料裡長出來的。"),

    ("CNN 在『看哪裡』:Grad-CAM",
     "打開黑盒的注意力",
     ["熱力圖疊在數字上,顯示判斷時注意力落在哪塊像素",
      "紅色=關鍵區;模型確實盯著筆畫,不是看背景 → 辨識更可信"],
     ["gradcam.png"],
     "這回答了『黑盒到底看哪裡』:它確實盯著筆畫。"),

    ("訓練與評估",
     "訓練曲線 · 混淆矩陣 · 特徵分群",
     ["訓練曲線:loss 下降、accuracy 上升 = 模型在學習",
      "混淆矩陣:看它最容易混淆哪些數字對",
      "t-SNE:倒數第二層特徵把 1~9 與空格分成 10 群 → 學到可分離的語意"],
     ["training_curve.png", "confusion_matrix.png", "tsne.png"],
     "t-SNE 證明學到的是可分離的表徵、不是死背;評估講得出道理比準確率更加分。"),

    ("⭐ Sim-to-Real:模型的真實極限",
     "乾淨近乎完美,唯一弱點是光線",
     ["乾淨數位輸入 ~100%、合成測試 97.7% — 幾乎無落差",
      "壓力測試:對焦模糊/感光雜訊/JPEG 都很穩",
      "唯一明顯弱點 = 光線不均(準確率明顯下滑)",
      "正好解釋拍照提示『光線均勻』;改進:光照校正 / 光照增強"],
     ["robustness_curves.png", "real_vs_synth.png"],
     "我沒硬湊落差——乾淨圖就是 100%;改用壓力測試找真正的極限,結果是光線。"),

    ("⭐ 資料增強消融(Ablation)",
     "用數字證明『增強』的價值",
     ["有增強 vs 無增強,同結構模型對照",
      "帶擾動合成測試:98.5% vs 70.0%(+28pp)",
      "真實照片:有增強更穩;劣化條件下差距更明顯",
      "消融 = 拿掉一個設計、其餘不變,證明它不是裝飾"],
     ["ablation_bar.png", "ablation_robustness.png"],
     "資料增強買到了實打實的穩健性,不是裝飾。"),

    ("⭐ 實戰除錯:白盒抓到黑盒的錯",
     "把空格幻覺『修在源頭』",
     ["真實照片把『空角落』誤判成數字 → 推理引擎回報『無解』",
      "合法數獨一定有解 → 白盒安全網確認:辨識錯了",
      "根因:空格殘留外框亮線,模型訓練時沒看過(空格都是純黑)",
      "修法:只給『空格類別』加邊框殘留增強 → 該情境 94.5%→~100%,乾淨仍 100%",
      "關鍵:增強只能加空格;連數字也加會讓模型被邊緣帶偏、讀不出淡數字"],
     ["grid_compare.png"],
     "三件事一次發生:白盒安全網作動、可複現的失敗、用針對性增強修在源頭。這才是理解 CNN。"),

    ("推理模組:白盒技巧引擎",
     "每步有名字、有理由、可驗證",
     ["四招階梯:唯一候選數★ → 隱性唯一數★★ → 顯性數對★★★ → 區塊摒除★★★☆",
      "每步附中文理由;難度評級 = 解開此題所需的『最難技巧』",
      "超出範圍 → 回溯法備援,並誠實標記"],
     ["confidence.png"],
     "這就是白盒:每步可驗證,和會看不會講的 CNN 形成對照。"),

    ("教練設計 + Live Demo",
     "拍照→辨識→自己填→要提示給下一步(不劇透)",
     ["主角:桌面 App play_gui.py — 拍照辨識、自己填、卡住按『下一步提示』",
      "提示從『目前盤面』即時算出 → 任何填字順序都對得上",
      "第一次用到某招附心法;一次只給一步,完整答案只在投降時出現",
      "填錯即時變紅警告(但不擋你填);備援:CLI play.py(零依賴)"],
     ["demo_mock.png"],
     "提示是從你當下盤面算出來的,所以順序無關;目標是讓玩家自己變強。"),

    ("限制與展望",
     "誠實面對",
     ["域偏移泛化不穩:合成→某張真實照片 run-to-run 會差(用 restore_best 挑最穩)",
      "尚未支援手寫數字;進階技巧(X-Wing 等)未實作",
      "工程取捨:辨識(重、需 TF)與遊玩(輕、零依賴)刻意解耦",
      "未來:鏡頭常開的 AR 即時版(只認變化的那格)、手機 App"],
     [],
     "會看的代價很高,所以讓 CNN 只在開局出場一次;之後提示是微秒級。"),

    ("結論",
     "把那顆 CNN 拆開看",
     ["學到邊緣與筆畫、注意力落在筆畫上、乾淨圖近乎完美",
      "唯一弱點是光線、增強值 +28pp、空格殘留會害它幻覺 — 而我修在了源頭",
      "有精確解的問題,別用機率的鎚子敲;用了機率工具,也要量得出它的邊界"],
     [],
     "神經網路會看、符號邏輯會想,讓它們各司其職。謝謝大家。"),
]


def add_textbox(slide, left, top, width, height, lines, size, color,
                bold=False, align=PP_ALIGN.LEFT, bullet=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = ("• " + line) if bullet else line
        run.font.size = Pt(size)
        run.font.name = FONT
        run.font.bold = bold
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return tb


def fit_pictures(slide, paths, top, band_h):
    """把多張圖等寬排在一條水平帶裡,維持比例置中"""
    from PIL import Image
    paths = [p for p in paths if p]
    if not paths:
        return
    n = len(paths)
    margin = Inches(0.4)
    gap = Inches(0.25)
    avail_w = EMU_W - 2 * margin - gap * (n - 1)
    cell_w = int(avail_w / n)
    x = margin
    for p in paths:
        iw, ih = Image.open(p).size
        scale = min(cell_w / iw, band_h / ih)
        w, h = int(iw * scale), int(ih * scale)
        off_x = x + (cell_w - w) // 2
        off_y = top + (band_h - h) // 2
        slide.shapes.add_picture(p, off_x, off_y, width=w, height=h)
        x += cell_w + gap


def main():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]

    for idx, (title, sub, bullets, imgs, notes) in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)

        if idx == 0:    # 標題頁置中大字
            add_textbox(slide, Inches(1), Inches(2.6), Inches(11.3), Inches(1.5),
                        [title], 40, NAVY, bold=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(1.2),
                        [sub], 18, GREY, align=PP_ALIGN.CENTER)
            slide.notes_slide.notes_text_frame.text = notes
            continue

        # 標題列
        add_textbox(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.9),
                    [title], 28, NAVY, bold=True)
        add_textbox(slide, Inches(0.55), Inches(1.05), Inches(12.3), Inches(0.5),
                    [sub], 15, ACCENT, bold=True)

        img_paths = [fig(n) for n in imgs]
        has_img = any(img_paths)

        if has_img and bullets:
            # 上方條列,下方圖帶
            add_textbox(slide, Inches(0.55), Inches(1.55), Inches(12.2), Inches(2.0),
                        bullets, 15, GREY, bullet=True)
            fit_pictures(slide, img_paths, Inches(3.7), Inches(3.5))
        elif has_img:
            fit_pictures(slide, img_paths, Inches(1.7), Inches(5.4))
        else:
            add_textbox(slide, Inches(0.7), Inches(1.8), Inches(12.0), Inches(4.8),
                        bullets, 20, GREY, bullet=True)

        slide.notes_slide.notes_text_frame.text = notes

    out = os.path.join(ROOT, "report", "SudokuCoach.pptx")
    prs.save(out)
    print(f"已產出簡報:{out}")
    print(f"共 {len(SLIDES)} 頁;口條已放進每頁的『簡報者備忘稿』")


if __name__ == "__main__":
    main()
